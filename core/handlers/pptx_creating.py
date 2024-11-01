import datetime
import os
from io import BytesIO

from PIL import Image

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
from pptx.dml.color import RGBColor
from moviepy.editor import VideoFileClip

MAX_WIDTH = Inches(13.33)
MAX_HEIGHT = Inches(7.5)


async def create_presentation(directory: str):
    """
    Создает презентации для каждой из команд в указанной директории

    :param directory: str путь к директории с файлами
    :return: list[str] список имен созданных презентаций
    """
    # Создание словаря для хранения файлов по командам
    teams_files = {}
    time = datetime.datetime.strftime(datetime.datetime.now(), '%H.%M.%S')
    # Получение списка файлов в директории
    async for file in get_files(directory):
        file_name = file.name
        # Разделение имени файла: тип файла, команда, уникальный код
        parts = file_name.split('_')
        if len(parts) < 3:
            continue  # Пропускаем файлы с некорректными именами
        team_name = time + '_' + parts[1]
        team_files = teams_files.get(team_name, [])
        team_files.append(file)
        teams_files[team_name] = team_files

    # Создание презентаций для каждой команды
    for team_name, files in teams_files.items():
        prs = Presentation()
        prs.slide_width = MAX_WIDTH  # 16:9 формат
        prs.slide_height = MAX_HEIGHT
        # создание слайда-заготовки
        slide_template = prs.slide_layouts[5]
        # установка черного фона на слайде
        background = slide_template.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        for file in files:
            file_extension = file.suffix.lower()

            slide = prs.slides.add_slide(slide_template)  # Пустой слайд
            if file_extension == '.jpg':
                # Добавляем изображение
                await add_image_to_slide(slide, file)
            elif file_extension in ('.mp4', '.mov'): # Добавляем видео
                # Выделяем первый кадр
                first_frame_pic = await first_frame_maker(file)
                first_frame_pic.seek(0)
                # создаем слайд с заглушкой перед слайдом с видео
                await add_static_slide(slide, file, first_frame_pic)
                # создаем сдайд с видео
                slide = prs.slides.add_slide(slide_template)  # Пустой слайд
                await add_video_to_slide(slide, file, first_frame_pic)

        # Сохранение презентации
        filepath = Path(directory, team_name).with_suffix('pptx').absolute().as_posix()
        prs.save(filepath)
    return teams_files.keys()


async def get_files(directory: str):
    """
    Генератор, возвращающий пути до файлов.
    """
    path = Path(directory)
    paths = list(path.iterdir())
    paths.sort(key=os.path.getctime)
    for file in paths:
        if file.is_file() and file.suffix.lower() in ['.jpg', '.mp4', '.mov']:
            yield file


async def coordinates_calculating(image_size) -> tuple:
    """
    Вычисляет координаты для масштабирования изображения,
    чтобы он был размещен по центру на слайде презентации
    и не выходил за границы слайда.

    :param image_size: размер изображения (ширина, высота)
    :return: кортеж с координатами левого верхнего угла и новыми размерами изображения
    """
    # Вычисление координат для масштабирования изображения
    width, height = image_size
    # Масштабируем изображение
    width_ratio = MAX_WIDTH / width
    height_ratio = MAX_HEIGHT / height
    scale_ratio = min(width_ratio, height_ratio)

    new_width = width * scale_ratio
    new_height = height * scale_ratio

    left = (MAX_WIDTH - new_width) / 2
    top = (MAX_HEIGHT - new_height) / 2

    return left, top, new_width, new_height


# Добавление изображения на слайд
async def add_image_to_slide(slide, image_path):
    """
    Функция добавления фото на слайд
    :param slide: слайд, на который добавляется изображение
    :param image_path: путь к изображению
    """
    im = Image.open(str(image_path))

    left, top, width, height = await coordinates_calculating(im.size)

    slide.shapes.add_picture(str(image_path),
                             left=left,
                             top=top,
                             width=width,
                             height=height)


# Добавление видео на слайд
async def add_video_to_slide(slide, video_path, fisrt_frame):
    """
    Добавление видео на слайд

    :param slide: слайд, на который добавляется видео
    :param video_path: путь к видеофайлу
    :return: None
    Превью к видео устнавливается первый кадр из видео.
    """
    clip = VideoFileClip(str(video_path))

    left, top, width, height = await coordinates_calculating(clip.size)

    slide.shapes.add_movie(str(video_path),
                           left=left,
                           top=top,
                           width=width,
                           height=height,
                           poster_frame_image=fisrt_frame,
                           )


async def add_static_slide(slide, video_path, fisrt_frame):
    """
    Добавление первого кадра видео на слайд, который будет идти до видео

    :param slide: слайд, на который добавляется видео
    :param video_path: путь к видеофайлу
    :first_frame: первый кадр из видео
    Превью к видео устнавливается первый кадр из видео.
    """
    clip = VideoFileClip(str(video_path))

    left, top, width, height = await coordinates_calculating(clip.size)

    slide.shapes.add_picture(fisrt_frame,
                             left=left,
                             top=top,
                             width=width,
                             height=height)


async def first_frame_maker(videofile_path: str) -> BytesIO:
    """
    Создает байтовый поток из первого кадра видео

    :param videofile_path: путь к видеофайлу
    :return BytesIO: байтовый поток из первого кадра видео
    """
    clip = VideoFileClip(str(videofile_path))

    first_frame = clip.get_frame(0)
    first_frame_image = Image.fromarray(first_frame)

    image_stream = BytesIO()
    first_frame_image.save(image_stream, format="PNG")
    return image_stream
